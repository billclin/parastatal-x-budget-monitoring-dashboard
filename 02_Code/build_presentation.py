"""
Parastatal X - 15-Slide Executive Presentation Builder
Uses python-pptx with forest green/muted gold theme
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

from pathlib import Path
BASE = Path(__file__).resolve().parents[1]

FGRN = RGBColor(0x1B, 0x43, 0x32)
GOLD = RGBColor(0xB7, 0x95, 0x0B)
CHRC = RGBColor(0x2C, 0x3E, 0x50)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
LGRY = RGBColor(0xF2, 0xF2, 0xF2)
LGRN = RGBColor(0xD5, 0xE8, 0xD4)
LGLD = RGBColor(0xFF, 0xF2, 0xCC)
RED  = RGBColor(0xC0, 0x39, 0x2B)
AMBR = RGBColor(0xE6, 0x7E, 0x22)
GRNS = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = 6  # blank slide layout index

def add_slide(prs):
    layout = prs.slide_layouts[BLANK]
    return prs.slides.add_slide(layout)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h, size=12, bold=False, color=CHRC,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=""):
    """Standard dark green header bar across top"""
    rect(slide, 0, 0, 13.33, 1.2, fill=FGRN)
    txb(slide, title, 0.3, 0.1, 9, 0.6, size=22, bold=True, color=WHT)
    if subtitle:
        txb(slide, subtitle, 0.3, 0.7, 10, 0.4, size=11, color=RGBColor(0xD4,0xE6,0xD6))
    # Gold accent line
    rect(slide, 0, 1.2, 13.33, 0.04, fill=GOLD)

def footer_bar(slide, text="Parastatal X | Financial Analytics Portfolio Project | All Data Synthetic | June 2026"):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=FGRN)
    txb(slide, text, 0.2, 7.12, 12, 0.3, size=8, color=WHT, italic=True)

def kpi_card(slide, l, t, w, h, label, value, unit="", bg=LGRY, val_color=FGRN):
    rect(slide, l, t, w, h, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, value, l+0.1, t+0.15, w-0.2, h*0.4, size=20, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txb(slide, label, l+0.05, t+h*0.5, w-0.1, h*0.4, size=8, color=CHRC, align=PP_ALIGN.CENTER)
    if unit:
        txb(slide, unit, l+0.05, t+h-0.3, w-0.1, 0.25, size=7, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
rect(slide, 0, 0, 13.33, 7.5, fill=FGRN)
rect(slide, 0, 5.5, 13.33, 2.0, fill=RGBColor(0x12,0x2D,0x22))
rect(slide, 0.5, 4.9, 12.33, 0.05, fill=GOLD)
txb(slide, "PARASTATAL X", 1.0, 0.8, 11.33, 1.2, size=42, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txb(slide, "DYNAMIC BUDGET MONITORING, FORECASTING", 0.5, 2.0, 12.33, 0.8, size=24, bold=True, color=WHT, align=PP_ALIGN.CENTER)
txb(slide, "AND MANAGEMENT DASHBOARD", 0.5, 2.7, 12.33, 0.7, size=24, bold=True, color=WHT, align=PP_ALIGN.CENTER)
txb(slide, "A Simulated Kenyan Public-Sector Financial Analytics Solution", 0.5, 3.5, 12.33, 0.5, size=14, color=RGBColor(0xD4,0xE6,0xD6), align=PP_ALIGN.CENTER, italic=True)
txb(slide, "Financial Year 2025/2026  |  Analysis Cut-off: 31 May 2026  |  Currency: Kenya Shillings (KSh)", 0.5, 4.1, 12.33, 0.4, size=11, color=RGBColor(0xD4,0xE6,0xD6), align=PP_ALIGN.CENTER)
txb(slide, "⚠  All data in this presentation is entirely synthetic. No real persons, institutions or government data are represented.", 0.5, 5.6, 12.33, 0.6, size=9, color=RGBColor(0xAA,0xCC,0xAA), align=PP_ALIGN.CENTER, italic=True)
txb(slide, "Tools: Python  |  Power BI  |  DAX  |  Power Query  |  Excel  |  XlsxWriter  |  DOCX  |  PPTX", 0.5, 6.3, 12.33, 0.4, size=9, color=RGBColor(0x88,0xBB,0x88), align=PP_ALIGN.CENTER)
txb(slide, "June 2026", 0.5, 6.9, 12.33, 0.4, size=10, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: BACKGROUND AND BUSINESS PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Background and Business Problem", "The financial management challenge in Kenyan parastatals")
rect(slide, 0.3, 1.4, 12.73, 5.5, fill=WHT)

txb(slide, "The Problem", 0.5, 1.5, 4.0, 0.4, size=13, bold=True, color=FGRN)
problems = [
    "📊  Budget data fragmented across multiple systems",
    "⏰  Reports delayed — weeks old by the time they reach management",
    "🚫  No systematic expenditure forecasting capability",
    "⚠️  Approval bottlenecks invisible until year-end",
    "💸  Commitment exposure hidden beneath the surface",
    "📋  Audit findings repeat year after year",
    "🔍  No drill-through from summary to transaction level",
]
for i, p in enumerate(problems):
    txb(slide, p, 0.5, 1.9+i*0.48, 6.2, 0.45, size=10, color=CHRC)

txb(slide, "Why This Matters", 7.0, 1.5, 6.0, 0.4, size=13, bold=True, color=FGRN)
context = [
    "🏛️  Kenyan PFM Act requires strong accountability",
    "📅  Jul–Jun financial year with year-end pressure",
    "🏢  32 cost centres across 13 directorates, 8 regions",
    "💰  KSh 1.70B budget requiring continuous monitoring",
    "🤝  Donor funds requiring separate tracking and reporting",
    "📈  Performance contracting demands evidence",
    "✅  Board and Treasury reporting expectations",
]
for i, c in enumerate(context):
    txb(slide, c, 7.0, 1.9+i*0.48, 6.0, 0.45, size=10, color=CHRC)

rect(slide, 0.3, 5.35, 12.73, 0.7, fill=LGLD)
txb(slide, '💡  Solution: A near-real-time financial monitoring dashboard consolidating all data into one view, enabling management to detect exceptions, track forecasts, and take informed action before year-end.', 0.5, 5.4, 12.3, 0.6, size=10, bold=True, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: PROJECT OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Project Objectives", "What this solution is designed to achieve")
objectives = [
    ("Budget Control","Provide near-real-time visibility of budget utilisation, variance, and absorption rates across all 32 cost centres and 13 directorates"),
    ("Expenditure Monitoring","Track all expenditure transactions, identify anomalies, monitor pending bills, and maintain a complete audit trail"),
    ("Commitment Management","Expose commitment exposure beyond YTD expenditure; identify overdue and high-risk commitments"),
    ("Approval Workflow","Monitor AIE request processing, detect bottlenecks, track aging requests, and measure turnaround time"),
    ("Forecasting","Apply 4-method blended year-end forecasting to provide management with advance warning of over-expenditure risk"),
    ("Revenue Tracking","Monitor AIA and revenue collection against targets; identify underperforming streams and regions"),
    ("Performance Integration","Link budget expenditure to service delivery outputs — enabling cost per output and value-for-money analysis"),
    ("Audit Readiness","Track audit exceptions, monitor resolution, flag repeat findings, and maintain a clean data quality log"),
]
for i, (title, detail) in enumerate(objectives):
    col = 0 if i < 4 else 6.5
    row = 1.35 + (i % 4) * 1.45
    rect(slide, col+0.2, row, 6.0, 1.35, fill=LGRN if i%2==0 else LGLD, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, f"{'●'} {title}", col+0.35, row+0.08, 5.6, 0.35, size=11, bold=True, color=FGRN)
    txb(slide, detail, col+0.35, row+0.43, 5.6, 0.85, size=9, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: SYNTHETIC DATASET OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Synthetic Dataset Overview", "18,896+ rows across 13 datasets — 3 financial years")
txb(slide, "⚠ All data is entirely synthetic. Random seed 42. No real data used.", 0.3, 1.3, 12.7, 0.35, size=9, italic=True, color=AMBR)

datasets = [
    ("Approved Annual Budget","96 rows","3 FYs × 32 cost centres","KSh 4.65B total"),
    ("Monthly Alloc.","1,152 rows","12 months × 32 CCs × 3 FYs","100% reconciled"),
    ("Expenditure Txns","12,100 rows","FY2023/24 to FY2025/26","KSh 4.13B total"),
    ("Commitments","1,350 rows","~450 per FY","KSh 700M open"),
    ("Budget Requests","1,000 rows","AIE workflow","KSh 2.48B backlog"),
    ("Pending Bills","520 rows","~173 per FY","KSh 1.12B outstanding"),
    ("Revenue/AIA","1,140 rows","8 revenue streams","88.2% achievement"),
    ("Procurement Plan","390 rows","~130 per FY","12 delayed contracts"),
    ("Performance","165 rows","15 KPIs × 3 FYs","~88% avg achievement"),
    ("Audit Exceptions","300 rows","~100 per FY","20% repeat findings"),
    ("Supplementary","209 rows","Budget adjustments","65% increases"),
    ("Reallocations","224 rows","75% approved","Net to zero"),
    ("Data Quality Log","250 rows","2.6% error rate","96.2% DQ score"),
]
cols_per_row = 4
for i, (name, rows, desc, kpi) in enumerate(datasets):
    col = (i % cols_per_row) * 3.2 + 0.15
    row = 1.7 + (i // cols_per_row) * 1.45
    bg = LGRN if i%3==0 else (LGLD if i%3==1 else LGRY)
    rect(slide, col, row, 3.1, 1.35, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, name, col+0.1, row+0.06, 2.9, 0.3, size=9, bold=True, color=FGRN)
    txb(slide, rows, col+0.1, row+0.36, 2.9, 0.25, size=14, bold=True, color=CHRC)
    txb(slide, desc, col+0.1, row+0.62, 2.9, 0.3, size=8, color=CHRC)
    txb(slide, kpi, col+0.1, row+0.95, 2.9, 0.3, size=8, bold=True, color=GOLD)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: DATA ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Data Architecture", "Star schema — 12 fact tables, 14 dimensions, 70+ DAX measures")

# Dim_Date centre
rect(slide, 5.7, 3.0, 2.0, 0.8, fill=FGRN)
txb(slide, "Dim_Date", 5.7, 3.1, 2.0, 0.6, size=11, bold=True, color=WHT, align=PP_ALIGN.CENTER)

facts = [
    (0.3, 1.4, "Fact_Approved_Budget"),
    (3.0, 1.4, "Fact_Monthly_Budget"),
    (6.0, 1.4, "Fact_Supplementary"),
    (9.0, 1.4, "Fact_Reallocation"),
    (0.3, 5.5, "Fact_Expenditure"),
    (3.0, 5.5, "Fact_Commitments"),
    (6.0, 5.5, "Fact_Pending_Bills"),
    (9.0, 5.5, "Fact_Revenue"),
    (11.5,3.2, "Fact_Procurement"),
    (11.5,4.2, "Fact_Performance"),
    (0.3, 3.2,  "Fact_Control_Exc."),
    (0.3, 4.2,  "Fact_DQ_Log"),
]
for fl, ft, fn in facts:
    rect(slide, fl, ft, 2.3, 0.7, fill=LGRN, line=RGBColor(0x1B,0x43,0x32))
    txb(slide, fn, fl+0.05, ft+0.15, 2.2, 0.45, size=8, bold=True, color=FGRN, align=PP_ALIGN.CENTER)

dims = [
    (4.5, 2.2, "Dim_Cost_Centre"),
    (4.5, 3.0, "Dim_Directorate"),
    (4.5, 3.8, "Dim_Region"),
    (4.5, 4.6, "Dim_Programme"),
    (7.8, 2.2, "Dim_Vote"),
    (7.8, 3.0, "Dim_Fund_Source"),
    (7.8, 3.8, "Dim_Activity"),
    (7.8, 4.6, "Dim_Donor"),
]
for dl, dt, dn in dims:
    rect(slide, dl, dt, 2.0, 0.65, fill=LGLD, line=RGBColor(0xB7,0x95,0x0B))
    txb(slide, dn, dl+0.05, dt+0.15, 1.9, 0.4, size=8, color=CHRC, align=PP_ALIGN.CENTER)

txb(slide, "Python → CSV → Power Query → Star Schema → DAX → Dashboard", 2.0, 6.9, 9.0, 0.4, size=9, color=WHT, italic=True, align=PP_ALIGN.CENTER)
rect(slide, 0, 6.8, 13.33, 0.7, fill=FGRN)
txb(slide, "Python → CSV → Power Query → Star Schema → DAX → Dashboard", 0.3, 6.85, 12.73, 0.4, size=9, color=WHT, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: BUDGET OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Budget Overview — FY2025/2026", "Approved KSh 1.70B | YTD Actual KSh 1.478B | Cut-off: 31 May 2026")

kpis = [
    ("Approved Budget", "KSh 1.700B", "Annual ceiling", LGRY, FGRN),
    ("YTD Expenditure", "KSh 1.478B", "As at 31 May 2026", LGLD, AMBR),
    ("Open Commitments", "KSh 0.700B", "Cash exposure", LGLD, AMBR),
    ("Available Balance", "KSh 0.222B", "Net of commits", LGLD, CHRC),
    ("Forecast YE", "KSh 1.763B", "103.7% absorption", LGLD, RED),
    ("Revenue Collected", "KSh 0.980B", "88.2% vs target", LGLD, AMBR),
]
for i, (lbl, val, sub, bg, vc) in enumerate(kpis):
    col = (i % 3) * 4.3 + 0.3
    row = 1.4 + (i // 3) * 1.8
    kpi_card(slide, col, row, 4.0, 1.6, sub, val, lbl, bg=bg, val_color=vc)

# Absorption bar
rect(slide, 0.5, 5.5, 12.33, 0.5, fill=LGRY)
rect(slide, 0.5, 5.5, 12.33*0.870, 0.5, fill=GRNS)
txb(slide, "YTD Absorption: 87.0%   |   Expected (time-based): 91.7%   |   Gap: -4.7 percentage points   →   AMBER", 0.5, 5.5, 12.33, 0.5, size=11, bold=True, color=WHT, align=PP_ALIGN.CENTER)
txb(slide, "Budget Profile:  ██████████████████████████████████████░░░░  87.0% of KSh 1.70B utilised with 11 of 12 months elapsed", 0.3, 6.1, 12.5, 0.4, size=9, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: EXPENDITURE PERFORMANCE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Expenditure Performance", "Monthly trend, directorate absorption, and year-on-year comparison")

txb(slide, "Monthly Expenditure Trend (FY2025/2026) — KSh Millions", 0.3, 1.3, 7.5, 0.35, size=11, bold=True, color=FGRN)
months = ["Jul","Aug","Sep","Oct","Nov","Dec","Jan","Feb","Mar","Apr","May"]
values = [98, 82, 118, 105, 112, 94, 108, 115, 128, 135, 159]
max_v = max(values)
for i, (m, v) in enumerate(zip(months, values)):
    bar_h = v / max_v * 3.5
    col = 0.4 + i * 1.15
    rect(slide, col, 5.1 - bar_h, 0.9, bar_h, fill=FGRN)
    txb(slide, m, col, 5.15, 0.9, 0.3, size=8, color=CHRC, align=PP_ALIGN.CENTER)
    txb(slide, str(v), col, 5.1 - bar_h - 0.35, 0.9, 0.32, size=8, bold=True, color=FGRN, align=PP_ALIGN.CENTER)

txb(slide, "Avg monthly spend: KSh 134.4M | June forecast: KSh 221.6M required (65% above avg)", 0.3, 5.55, 12.5, 0.3, size=9, italic=True, color=CHRC)

txb(slide, "Absorption by Directorate", 8.3, 1.3, 4.8, 0.35, size=11, bold=True, color=FGRN)
dirs = [
    ("ICT Directorate","93.8%","Green"),
    ("Finance & Accounts","92.1%","Green"),
    ("Human Resources","91.5%","Green"),
    ("Strategy & Planning","85.2%","Amber"),
    ("Procurement","79.8%","Amber"),
    ("Projects & Infra","74.1%","Red"),
    ("Regional Operations","71.3%","Red"),
    ("Research & M&E","68.9%","Red"),
]
for i, (d, pct, rag) in enumerate(dirs):
    row = 1.75 + i * 0.46
    bg = LGRN if rag=="Green" else (LGLD if rag=="Amber" else RGBColor(0xFA,0xDB,0xD8))
    vc = GRNS if rag=="Green" else (AMBR if rag=="Amber" else RED)
    rect(slide, 8.3, row, 4.8, 0.42, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, d, 8.4, row+0.06, 3.2, 0.32, size=8.5, color=CHRC)
    txb(slide, pct, 11.5, row+0.06, 1.4, 0.32, size=9, bold=True, color=vc, align=PP_ALIGN.RIGHT)

rect(slide, 0.3, 5.8, 12.73, 0.8, fill=LGLD)
txb(slide, "⚠  Year-on-Year Context:  FY2024/2025 total expenditure was KSh 1.338B. Current YTD of KSh 1.478B represents 10.3% growth — consistent with the 9.7% budget increase. YoY trajectory is positive; absorption rate is the primary concern.", 0.5, 5.82, 12.3, 0.74, size=9, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: APPROVAL AND PROCESSING WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Budget Review and Approval Workflow", "AIE processing, aging analysis, and bottleneck identification")
stages = [
    ("Received","1,000","KSh 3.25B"),
    ("Approved","~700","KSh 3.25B"),
    ("Processed","~350","KSh 0.77B"),
    ("Outstanding","~350","KSh 2.48B"),
]
for i, (stage, cnt, amt) in enumerate(stages):
    w = 3.0 - i * 0.3
    col = 0.3 + i * 3.1
    rect(slide, col, 1.4, w, 1.8, fill=FGRN if i==0 else (LGRN if i==1 else (LGLD if i==2 else RGBColor(0xFA,0xDB,0xD8))))
    vc = WHT if i==0 else FGRN if i==1 else AMBR if i==2 else RED
    txb(slide, stage, col, 1.5, w, 0.4, size=13, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txb(slide, cnt, col, 1.9, w, 0.45, size=18, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txb(slide, amt, col, 2.35, w, 0.4, size=10, color=vc, align=PP_ALIGN.CENTER)

txb(slide, "→", 3.3, 2.05, 0.8, 0.5, size=20, bold=True, color=FGRN, align=PP_ALIGN.CENTER)
txb(slide, "→", 6.4, 2.05, 0.8, 0.5, size=20, bold=True, color=AMBR, align=PP_ALIGN.CENTER)
txb(slide, "→", 9.5, 2.05, 0.8, 0.5, size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)

txb(slide, "Aging Analysis — FY2025/2026", 0.3, 3.35, 6.0, 0.35, size=11, bold=True, color=FGRN)
aging = [("0–14 Days","~28%","Green"),("15–30 Days","~22%","Amber"),("31–60 Days","~28%","Red"),("61–90 Days","~14%","Red"),("Over 90 Days","~8%","Red")]
for i, (band, pct, rag) in enumerate(aging):
    bg = LGRN if rag=="Green" else (LGLD if rag=="Amber" else RGBColor(0xFA,0xDB,0xD8))
    rect(slide, 0.3+i*2.42, 3.75, 2.3, 0.9, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    vc = GRNS if rag=="Green" else AMBR if rag=="Amber" else RED
    txb(slide, pct, 0.3+i*2.42, 3.8, 2.3, 0.4, size=16, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txb(slide, band, 0.3+i*2.42, 4.2, 2.3, 0.4, size=8.5, color=CHRC, align=PP_ALIGN.CENTER)

txb(slide, "Key Bottlenecks Identified", 7.0, 3.35, 6.0, 0.35, size=11, bold=True, color=FGRN)
bottlenecks = [
    "⚠  ~50% of requests pending >30 days (target: 14 days)",
    "⚠  KSh 2.48B approved but not yet processed",
    "⚠  Average approval turnaround: 18+ days",
    "⚠  Processing Rate: ~35% (target: >80%)",
    "✓  Approval Rate: ~70% (acceptable range)",
]
for i, b in enumerate(bottlenecks):
    txb(slide, b, 7.0, 3.75+i*0.46, 6.0, 0.42, size=9.5, color=CHRC)

rect(slide, 0.3, 5.6, 12.73, 0.7, fill=LGLD)
txb(slide, "Management Action: Deploy additional processing staff; implement 5-day SLA for all AIE requests; escalate stale requests weekly to Director, Finance and Accounts.", 0.5, 5.65, 12.3, 0.6, size=9.5, bold=True, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: COMMITMENTS AND PENDING BILLS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Commitments and Pending Bills", "Cash exposure beyond YTD expenditure — KSh 1.82B combined")

txb(slide, "OPEN COMMITMENTS", 0.3, 1.35, 6.2, 0.4, size=13, bold=True, color=FGRN)
commit_metrics = [
    ("Total Open Commitments", "KSh 700M", LGLD, AMBR),
    ("Commitments Overdue", "~15%", LGLD, RED),
    ("Avg Days Outstanding", "68 days", LGLD, AMBR),
    ("High Risk Commitments", "~12%", RGBColor(0xFA,0xDB,0xD8), RED),
]
for i, (lbl, val, bg, vc) in enumerate(commit_metrics):
    col = 0.3 + (i%2) * 3.0
    row = 1.8 + (i//2) * 1.05
    kpi_card(slide, col, row, 2.8, 0.95, lbl, val, "", bg=bg, val_color=vc)

txb(slide, "Commitment by Status", 0.3, 3.95, 6.2, 0.35, size=10, bold=True, color=FGRN)
statuses = [("Open","38%",RED),("Part. Invoiced","22%",AMBR),("Part. Paid","18%",GOLD),("Overdue","12%",RED),("Other","10%",LGRN)]
for i,(s,p,c) in enumerate(statuses):
    rect(slide, 0.3, 4.35+i*0.42, float(p.replace('%',''))/100*6, 0.38, fill=c)
    txb(slide, f"{s}  {p}", 0.4, 4.38+i*0.42, 5.8, 0.34, size=9, color=WHT, bold=True)

txb(slide, "PENDING BILLS", 6.8, 1.35, 6.2, 0.4, size=13, bold=True, color=FGRN)
pb_metrics = [
    ("Total Outstanding", "KSh 1.122B", RGBColor(0xFA,0xDB,0xD8), RED),
    ("Bills Over 90 Days", "~18%", RGBColor(0xFA,0xDB,0xD8), RED),
    ("Avg Bill Age", "47 days", LGLD, AMBR),
    ("Incomplete Docs", "~28%", LGLD, AMBR),
]
for i, (lbl, val, bg, vc) in enumerate(pb_metrics):
    col = 6.8 + (i%2) * 3.0
    row = 1.8 + (i//2) * 1.05
    kpi_card(slide, col, row, 2.8, 0.95, lbl, val, "", bg=bg, val_color=vc)

txb(slide, "Bills by Aging Band", 6.8, 3.95, 6.2, 0.35, size=10, bold=True, color=FGRN)
aging_pb = [("0–30 Days","45%",GRNS),("31–90 Days","37%",AMBR),("Over 90 Days","18%",RED)]
for i,(s,p,c) in enumerate(aging_pb):
    rect(slide, 6.8, 4.35+i*0.6, float(p.replace('%',''))/100*6, 0.52, fill=c)
    txb(slide, f"{s}: {p}", 6.9, 4.38+i*0.6, 5.8, 0.46, size=9.5, color=WHT, bold=True)

rect(slide, 0.3, 6.05, 12.73, 0.65, fill=RGBColor(0xFA,0xDB,0xD8))
txb(slide, "🚨  Combined Exposure: KSh 1.822B (KSh 700M commitments + KSh 1.122B pending bills) — equivalent to 107% of approved annual budget. Immediate management action required.", 0.5, 6.07, 12.3, 0.60, size=9.5, bold=True, color=RED)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: FORECAST AND RISK ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Forecast and Risk Analysis", "4-method blended forecast | 3-scenario analysis | RAG risk heatmap")

txb(slide, "Year-End Expenditure Forecast Methods", 0.3, 1.35, 7.5, 0.35, size=11, bold=True, color=FGRN)
methods = [
    ("Straight-Line Run Rate","KSh 1,612,762,691","94.9%","Conservative", LGRN, GRNS),
    ("3-Month Moving Average","KSh 1,653,229,667","97.2%","Moderate",LGLD, AMBR),
    ("Seasonal Adjustment","KSh 1,722,723,783","101.3%","Moderate",LGLD, AMBR),
    ("Commitment-Adjusted","KSh 2,178,365,800","128.1%","Extreme",RGBColor(0xFA,0xDB,0xD8), RED),
    ("BLENDED SELECTED","KSh 1,763,490,330","103.7%","Slight Over-exp",RGBColor(0xFA,0xDB,0xD8), RED),
]
for i,(m,f,pct,desc,bg,vc) in enumerate(methods):
    is_selected = i == 4
    rect(slide, 0.3, 1.75+i*0.83, 7.5, 0.78,
         fill=RGBColor(0xEB,0xF5,0xEB) if is_selected else bg,
         line=GOLD if is_selected else RGBColor(0xCC,0xCC,0xCC))
    txb(slide, m, 0.4, 1.8+i*0.83, 3.5, 0.35, size=9, bold=is_selected, color=FGRN if is_selected else CHRC)
    txb(slide, f, 3.9, 1.8+i*0.83, 2.7, 0.35, size=10, bold=is_selected, color=vc, align=PP_ALIGN.RIGHT)
    txb(slide, pct, 6.65, 1.8+i*0.83, 1.0, 0.35, size=9, bold=is_selected, color=vc, align=PP_ALIGN.RIGHT)

txb(slide, "Scenario Analysis", 8.2, 1.35, 4.9, 0.35, size=11, bold=True, color=FGRN)
scenarios = [
    ("Base Case","KSh 1.763B","103.7%","KSh -63.5M",LGLD,AMBR),
    ("Optimistic","KSh 1.711B","100.6%","KSh -10.6M",LGLD,AMBR),
    ("Adverse","KSh 1.905B","112.0%","KSh -204.6M",RGBColor(0xFA,0xDB,0xD8),RED),
]
for i,(sc,fe,fp,bal,bg,vc) in enumerate(scenarios):
    rect(slide, 8.2, 1.75+i*1.5, 4.9, 1.4, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, sc, 8.3, 1.8+i*1.5, 4.7, 0.35, size=11, bold=True, color=vc)
    txb(slide, f"Forecast: {fe}  ({fp})", 8.3, 2.15+i*1.5, 4.7, 0.3, size=9, color=CHRC)
    txb(slide, f"Balance: {bal}", 8.3, 2.45+i*1.5, 4.7, 0.3, size=9, bold=True, color=RED)

rect(slide, 0.3, 6.15, 12.73, 0.55, fill=RGBColor(0xFA,0xDB,0xD8))
txb(slide, "⚠  Over-Expenditure Risk: KSh 63.5M. Blended weights: SL 30% | MA3 25% | Seasonal 25% | Commitment 20%. Required June spend to absorb budget: KSh 221.6M vs avg KSh 134.4M. Intervention required.", 0.5, 6.17, 12.3, 0.50, size=9, bold=True, color=RED)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: PERFORMANCE AND VALUE FOR MONEY
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Performance and Value for Money", "KPI achievement vs budget expenditure | 4-quadrant analysis | Cost per output")

txb(slide, "4-Quadrant Performance Analysis", 0.3, 1.35, 7.5, 0.35, size=11, bold=True, color=FGRN)
rect(slide, 0.3, 1.75, 6.0, 4.5, fill=WHT, line=RGBColor(0x88,0x88,0x88))
rect(slide, 0.3, 1.75, 3.0, 2.25, fill=LGLD)  # Q2: low exp, high perf
rect(slide, 3.3, 1.75, 3.0, 2.25, fill=LGRN)  # Q1: high exp, high perf
rect(slide, 0.3, 4.0, 3.0, 2.25, fill=RGBColor(0xFA,0xDB,0xD8))  # Q3: low exp, low perf
rect(slide, 3.3, 4.0, 3.0, 2.25, fill=LGLD)  # Q4: high exp, low perf
txb(slide, "Low Spend\nHigh Performance\n✓ EFFICIENT", 0.5, 2.0, 2.7, 1.8, size=9, color=CHRC, align=PP_ALIGN.CENTER)
txb(slide, "High Spend\nHigh Performance\n✓ EFFECTIVE", 3.5, 2.0, 2.7, 1.8, size=9, color=GRNS, bold=True, align=PP_ALIGN.CENTER)
txb(slide, "Low Spend\nLow Performance\n✗ AT RISK", 0.5, 4.25, 2.7, 1.7, size=9, color=RED, bold=True, align=PP_ALIGN.CENTER)
txb(slide, "High Spend\nLow Performance\n✗ POOR VFM", 3.5, 4.25, 2.7, 1.7, size=9, color=AMBR, bold=True, align=PP_ALIGN.CENTER)
txb(slide, "← Budget Expenditure →", 0.3, 6.3, 6.0, 0.3, size=8, color=CHRC, italic=True, align=PP_ALIGN.CENTER)
txb(slide, "↑ KPI Achievement ↑", 0.0, 3.5, 0.35, 2.0, size=7, color=CHRC, italic=True)

txb(slide, "KPI Performance Summary", 6.8, 1.35, 6.2, 0.35, size=11, bold=True, color=FGRN)
kpi_perf = [
    ("Infrastructure Projects","Annual Target: 8","Actual: 6","75.0%","Amber"),
    ("Budget Absorption Rate","Annual Target: 95%","Actual: 87%","91.6%","Amber"),
    ("Staff Trained","Annual Target: 80","Actual: 74","92.5%","Green"),
    ("Revenue Achievement","Annual Target: 95%","Actual: 88%","92.6%","Green"),
    ("Procurement Plan Exec.","Annual Target: 85%","Actual: 72%","84.7%","Amber"),
    ("M&E Reports Submitted","Annual Target: 4","Actual: 3","75.0%","Red"),
]
for i,(ind,tgt,act,pct,rag) in enumerate(kpi_perf):
    bg = LGRN if rag=="Green" else LGLD if rag=="Amber" else RGBColor(0xFA,0xDB,0xD8)
    vc = GRNS if rag=="Green" else AMBR if rag=="Amber" else RED
    rect(slide, 6.8, 1.75+i*0.85, 6.2, 0.78, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, ind, 6.9, 1.8+i*0.85, 3.5, 0.3, size=9, bold=True, color=FGRN)
    txb(slide, f"{tgt}  |  {act}", 6.9, 2.1+i*0.85, 3.8, 0.25, size=8, color=CHRC)
    txb(slide, pct, 10.7, 1.82+i*0.85, 2.0, 0.4, size=14, bold=True, color=vc, align=PP_ALIGN.RIGHT)

rect(slide, 0.3, 6.6, 12.73, 0.7, fill=LGLD)
txb(slide, "Average KPI Achievement: 88.3%  |  Financial Absorption: 87.0%  |  Financial–Physical Gap: +1.3 pts (financial slightly lags physical)  |  Average Cost per Output: Directorate-specific", 0.5, 6.62, 12.3, 0.65, size=9, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: DATA QUALITY AND CONTROLS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Data Quality and Audit Controls", "96.2% data quality score | 300 audit exceptions | Control monitoring")

txb(slide, "DATA QUALITY SCORE: 96.2%", 0.3, 1.35, 6.0, 0.5, size=14, bold=True, color=GRNS)
txb(slide, "Status: GREEN (above 95% amber threshold)", 0.3, 1.8, 6.0, 0.35, size=10, color=CHRC, italic=True)

dq_issues = [
    ("Issue Type","Count","Severity","Status"),
    ("Duplicate Transaction Ref","14","Critical","Resolved"),
    ("Processed > Approved Amount","8","Critical","Resolved"),
    ("Negative Expenditure Amounts","6","Critical","Resolved"),
    ("Missing Cost Centre Code","19","High","Resolved"),
    ("Incorrect Vote Code","23","High","Resolved"),
    ("Invalid Transaction Date","11","High","Resolved"),
    ("Inconsistent Region Names","31","Medium","Resolved"),
    ("Leading/Trailing Spaces","45","Low","Resolved"),
    ("Different Capitalisation","38","Low","Resolved"),
    ("Blank Approval Fields","27","Medium","In Review"),
    ("Formula/Type Errors","17","High","In Review"),
    ("Unmatched Lookup Codes","11","High","Pending"),
]
for i, row in enumerate(dq_issues):
    bg = FGRN if i==0 else (LGRD if i>0 and row[2]=="Critical" else (LGLD if row[2]=="High" else LGRY))
    LGRD = RGBColor(0xFA,0xDB,0xD8)
    if i == 0:
        for j, v in enumerate(row):
            rect(slide, 0.3+j*1.45, 2.2, 1.42, 0.4, fill=FGRN)
            txb(slide, v, 0.3+j*1.45, 2.23, 1.42, 0.36, size=8.5, bold=True, color=WHT, align=PP_ALIGN.CENTER)
    else:
        sev = row[2]; stat = row[3]
        bg = RGBColor(0xFA,0xDB,0xD8) if sev=="Critical" else (LGLD if sev=="High" else (LGRY if i%2==0 else WHT))
        vc_s = GRNS if stat=="Resolved" else AMBR if stat=="In Review" else RED
        for j, v in enumerate(row):
            rect(slide, 0.3+j*1.45, 2.2+i*0.43, 1.42, 0.41, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
            fc = vc_s if j==3 else CHRC
            txb(slide, v, 0.3+j*1.45, 2.24+i*0.43, 1.42, 0.36, size=8, color=fc, bold=(j==3 and stat=="Resolved"), align=PP_ALIGN.CENTER if j in [1,3] else PP_ALIGN.LEFT)

txb(slide, "AUDIT EXCEPTIONS SUMMARY", 6.4, 1.35, 6.7, 0.4, size=12, bold=True, color=FGRN)
exc_metrics = [
    ("Total Exceptions","300",LGRY,CHRC),
    ("Open/In Progress","~55%",LGLD,AMBR),
    ("Resolved","~30%",LGRN,GRNS),
    ("Repeat Findings","~20%",RGBColor(0xFA,0xDB,0xD8),RED),
    ("Amount at Risk","~KSh 45M",LGLD,RED),
]
for i,(lbl,val,bg,vc) in enumerate(exc_metrics):
    kpi_card(slide, 6.4+(i%2)*3.2, 1.8+(i//2)*1.2, 3.0, 1.1, lbl, val, "", bg=bg, val_color=vc)

txb(slide, "Top Exception Types:", 6.4, 4.2, 6.7, 0.35, size=10, bold=True, color=FGRN)
top_exc = ["Missing Supporting Documents","Incorrect Vote Code","Duplicate Payment Risk","Expenditure Above Budget","Long-Outstanding Commitment"]
for i,e in enumerate(top_exc):
    txb(slide, f"{'🔴' if i<2 else '🟡'}  {e}", 6.4, 4.58+i*0.42, 6.7, 0.4, size=9, color=CHRC)

rect(slide, 0.3, 6.6, 12.73, 0.65, fill=LGRN)
txb(slide, "✅  Reconciliation: All 18 reconciliation checks pass. Monthly allocations reconcile 100% to annual budgets. All pending bill outstanding amounts balance invoice minus paid.", 0.5, 6.62, 12.3, 0.60, size=9.5, color=FGRN, bold=True)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: MANAGEMENT RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Management Recommendations", "10 immediate actions for the Chief Executive Officer and Finance Director")

recs = [
    ("CRITICAL","CEO","Instruct bi-weekly monitoring of 23 Red cost centres; impose expenditure ceiling on 3 at-risk vote lines","By 10 Jun 2026"),
    ("CRITICAL","Dir. Finance","Review all open commitments above KSh 5M; cancel stale POs; prepare payment schedule","By 15 Jun 2026"),
    ("HIGH","Dir. Finance","Convene emergency year-end expenditure meeting across all directorates","By 10 Jun 2026"),
    ("HIGH","Head of Budget","Clear KSh 2.48B AIE backlog; implement 5-day SLA; escalate weekly","By 20 Jun 2026"),
    ("HIGH","Chief Accountant","Clear all pending bills >90 days; document payment schedule for Board","By 30 Jun 2026"),
    ("HIGH","Revenue Mgr","Launch June 2026 revenue drive; issue demand notices to debtors","By 10 Jun 2026"),
    ("HIGH","Project Manager","Submit revised donor fund implementation plan; notify donors","By 15 Jun 2026"),
    ("MEDIUM","Dir. Supply Chain","Issue contract extensions for 12 delayed contracts; re-procure where needed","By 30 Jun 2026"),
    ("MEDIUM","Head IA","Report on status of all critical audit exceptions; complete document drive","By 20 Jun 2026"),
    ("MEDIUM","Data Officer","Resolve all 35 critical/high DQ exceptions; update DQ log","By 15 Jun 2026"),
]
col_widths = [1.0, 1.8, 7.8, 2.0]
hdrs = ["Priority","Responsible","Recommended Action","Target Date"]
for c,(h,w) in enumerate(zip(hdrs,col_widths)):
    rect(slide, 0.15+sum(col_widths[:c]), 1.35, w, 0.45, fill=FGRN)
    txb(slide, h, 0.2+sum(col_widths[:c]), 1.37, w-0.1, 0.41, size=9, bold=True, color=WHT, align=PP_ALIGN.CENTER)
for i,(p,r,a,d) in enumerate(recs):
    bg = RGBColor(0xFA,0xDB,0xD8) if p=="CRITICAL" else (LGLD if p=="HIGH" else LGRY)
    pc = RED if p=="CRITICAL" else AMBR if p=="HIGH" else CHRC
    for c,(v,w) in enumerate(zip([p,r,a,d],col_widths)):
        rect(slide, 0.15+sum(col_widths[:c]), 1.8+i*0.52, w, 0.5, fill=bg if c==0 else (LGRY if i%2==0 else WHT), line=RGBColor(0xCC,0xCC,0xCC))
        fc = pc if c==0 else CHRC
        txb(slide, v, 0.2+sum(col_widths[:c]), 1.83+i*0.52, w-0.1, 0.45, size=8.5, bold=(c==0), color=fc, align=PP_ALIGN.CENTER if c in [0,3] else PP_ALIGN.LEFT)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: EXPECTED BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
header_bar(slide, "Expected Benefits", "Demonstrated and potential value of the analytics solution")
benefits = [
    ("Improved Budget Visibility","Consolidates 13 datasets into one view. Management sees approved budget, YTD spend, commitments, pending bills and forecasts on a single screen — replacing multiple disconnected spreadsheets."),
    ("Faster Exception Detection","RAG traffic lights flag 23 Red cost centres and 5 high-risk vote lines in seconds. Management can intervene months before year-end rather than discovering problems during audit."),
    ("Stronger Expenditure Control","Available-balance-after-commitments shows KSh 700M in hidden exposure not visible in simple budget-minus-expenditure reports. More conservative and accurate than conventional reporting."),
    ("Better Approval Monitoring","AIE workflow page reveals KSh 2.48B processing backlog — previously invisible in conventional reports. Aging analysis targets bottlenecks with precision."),
    ("Improved Forecasting","4-method blended forecast provides early warning of KSh 63.5M over-expenditure risk — actionable 6 weeks before year-end."),
    ("Value-for-Money Analysis","Links budget to service delivery outputs. Cost per output and 4-quadrant analysis identify where the organisation is efficient, effective, at-risk or poor value."),
    ("Stronger Audit Readiness","Exceptions dashboard tracks all 300 findings, monitors resolution, flags 20% repeat findings — creating systematic management response where none existed before."),
    ("Better Accountability","Management action register provides live tracking of 10 priority actions with responsible roles, target dates, progress percentages and follow-up comments."),
]
for i,(title,desc) in enumerate(benefits):
    col = (i%2) * 6.5 + 0.3
    row = 1.4 + (i//2) * 1.45
    bg = LGRN if i%4 in [0,3] else LGLD
    rect(slide, col, row, 6.2, 1.38, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
    txb(slide, f"{'✅'} {title}", col+0.15, row+0.07, 5.9, 0.36, size=10, bold=True, color=FGRN)
    txb(slide, desc, col+0.15, row+0.45, 5.9, 0.87, size=8.5, color=CHRC)
footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
rect(slide, 0, 0, 13.33, 7.5, fill=FGRN)
rect(slide, 0, 6.5, 13.33, 1.0, fill=RGBColor(0x12,0x2D,0x22))
rect(slide, 0.5, 1.5, 12.33, 0.06, fill=GOLD)
rect(slide, 0.5, 5.8, 12.33, 0.06, fill=GOLD)
txb(slide, "PARASTATAL X", 0.5, 0.3, 12.33, 0.9, size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txb(slide, "Dynamic Budget Monitoring, Forecasting and Management Dashboard", 0.5, 1.1, 12.33, 0.5, size=16, color=WHT, align=PP_ALIGN.CENTER)

summary = [
    "18,896+ rows of realistic synthetic data across 3 financial years",
    "32 cost centres | 13 directorates | 8 regions | 40 activities | 33 vote codes",
    "Complete data lifecycle: Python generation → cleaning → star schema → DAX → dashboard",
    "4-method blended forecast model | 3-scenario analysis | 36 validation checks",
    "70+ DAX measures | 500 lines of Power Query M code | Custom Power BI theme",
    "12-page dashboard design covering budget, expenditure, workflow, risk, performance",
    "Senior Finance Manager commentary | Action plan | Risk register | Case study",
]
for i, s in enumerate(summary):
    txb(slide, f"  ●  {s}", 1.5, 1.7+i*0.58, 10.33, 0.54, size=11, color=RGBColor(0xD4,0xE6,0xD6))

txb(slide, "This project demonstrates the full analytics stack applied to the Kenyan public sector:", 1.0, 5.9, 11.33, 0.4, size=9.5, color=RGBColor(0xD4,0xE6,0xD6), italic=True, align=PP_ALIGN.CENTER)
tools = "Python  •  Power BI  •  DAX  •  Power Query  •  Excel  •  Star Schema  •  Financial Forecasting  •  Management Reporting"
txb(slide, tools, 0.5, 6.3, 12.33, 0.5, size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txb(slide, "⚠ All data is entirely synthetic | Not representative of any real institution | June 2026", 0.5, 6.85, 12.33, 0.4, size=8, color=RGBColor(0xAA,0xCC,0xAA), italic=True, align=PP_ALIGN.CENTER)

filepath = f"{BASE}/06_Presentation/Parastatal_X_Budget_Dashboard_Presentation.pptx"
prs.save(filepath)
size = os.path.getsize(filepath) / (1024*1024)
print(f"Presentation saved: {filepath}")
print(f"Slides: 15  |  File size: {size:.1f} MB")
